import logging
import tempfile
from pathlib import Path

from ..config import (
    DASHSCOPE_API_KEY, DASHSCOPE_BASE_URL,
    EMBEDDING_MODEL, KNOWLEDGE_CHUNK_SIZE, KNOWLEDGE_CHUNK_OVERLAP,
    KNOWLEDGE_SEARCH_TOP_K, KNOWLEDGE_SEARCH_THRESHOLD,
)
from ..database import save_chunks, search_chunks
from openai import AsyncOpenAI

logger = logging.getLogger(__name__)


def _get_embedding_client() -> AsyncOpenAI | None:
    if not DASHSCOPE_API_KEY:
        return None
    return AsyncOpenAI(api_key=DASHSCOPE_API_KEY, base_url=DASHSCOPE_BASE_URL)


async def embed_text(texts: list[str]) -> list[list[float]]:
    client = _get_embedding_client()
    if not client:
        raise RuntimeError("DASHSCOPE_API_KEY 未配置")
    resp = await client.embeddings.create(model=EMBEDDING_MODEL, input=texts)
    embeddings = [item.embedding for item in resp.data]
    logger.info("向量化完成 | texts=%d dim=%d", len(texts), len(embeddings[0]) if embeddings else 0)
    return embeddings


def parse_pdf(file_path: str) -> str:
    import fitz
    doc = fitz.open(file_path)
    text = "\n".join(page.get_text() for page in doc)
    doc.close()
    return text


def parse_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)


def parse_xlsx(file_path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    texts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                texts.append(" | ".join(cells))
    wb.close()
    return "\n".join(texts)


_PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".doc": parse_docx,
    ".pptx": parse_pptx,
    ".ppt": parse_pptx,
    ".xlsx": parse_xlsx,
    ".xls": parse_xlsx,
}


def parse_document(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    parser = _PARSERS.get(ext)
    if not parser:
        raise ValueError(f"不支持的文件格式: {ext}")
    return parser(file_path)


def chunk_text(text: str, chunk_size: int = None, overlap: int = None) -> list[str]:
    if chunk_size is None:
        chunk_size = KNOWLEDGE_CHUNK_SIZE
    if overlap is None:
        overlap = KNOWLEDGE_CHUNK_OVERLAP
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    chunks = []
    current = []
    current_len = 0
    for para in paragraphs:
        if current_len + len(para) > chunk_size and current:
            chunks.append("\n".join(current))
            overlap_texts = current
            back_len = 0
            for t in reversed(current):
                if back_len + len(t) >= overlap:
                    break
                back_len += len(t)
            current = current[-len(overlap_texts):] if overlap_texts else []
            current = [t for t in overlap_texts if t][-1:] if current else []
            current_len = sum(len(t) for t in current)
        current.append(para)
        current_len += len(para)
    if current:
        chunks.append("\n".join(current))
    logger.info("文本切片完成 | origin_len=%d chunks=%d", len(text), len(chunks))
    return chunks


async def index_document(file_path: str, doc_id: str):
    logger.info("开始索引文档 | doc_id=%s path=%s", doc_id, file_path)
    text = parse_document(file_path)
    if not text.strip():
        logger.warning("文档内容为空 | doc_id=%s", doc_id)
        return
    chunks = chunk_text(text)
    embeddings = await embed_text(chunks)
    save_chunks(doc_id, chunks, embeddings)
    logger.info("文档索引完成 | doc_id=%s chunks=%d", doc_id, len(chunks))


async def search_knowledge(query: str, top_k: int = None) -> str:
    if top_k is None:
        top_k = KNOWLEDGE_SEARCH_TOP_K
    if not query.strip():
        return ""
    logger.info("知识库检索 | query=%s", query[:50])
    query_emb = await embed_text([query])
    results = search_chunks(query_emb[0], top_k)
    if not results:
        logger.info("知识库无匹配结果")
        return ""
    filtered = [r for r in results if r["score"] >= KNOWLEDGE_SEARCH_THRESHOLD]
    if not filtered:
        logger.info("所有结果均低于阈值 threshold=%s max_score=%s", KNOWLEDGE_SEARCH_THRESHOLD, results[0]["score"])
        return ""
    lines = []
    for i, r in enumerate(filtered, 1):
        lines.append(f"[{i}] (相关度: {r['score']:.2f}) {r['content']}")
    context = "以下是企业知识库中的相关内容，请优先参考这些信息来回答用户问题：\n\n"
    context += "\n\n".join(lines)
    context += "\n\n请结合知识库内容和你的知识来回答。如果在知识库中找到了相关信息，请优先采用知识库的内容。"
    logger.info("知识库检索结果 | matched=%d top_score=%s", len(filtered), filtered[0]["score"])
    return context
